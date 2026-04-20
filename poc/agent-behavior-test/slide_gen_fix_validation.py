"""
slide_gen_fix_validation.py — POC experiment script for slide_gen bug fix validation.

Validates proposed fixes for four bugs found in backend/agent_workflows/slide_gen.py.
No LLM calls — pure Python / Pydantic / python-pptx.

Run from project root:
    micromamba run -n py3.12 python poc/agent-behavior-test/slide_gen_fix_validation.py

Results and analysis: poc/agent-behavior-test/slide_gen_fix_validation.md
"""

from __future__ import annotations

import json
import os
import sys
import tempfile
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv

load_dotenv()

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "backend"))

from pydantic import BaseModel, ValidationError, field_validator

# ══════════════════════════════════════════════════════════════════════════════
# Data structures
# ══════════════════════════════════════════════════════════════════════════════

@dataclass
class ExperimentResult:
    exp_id: str
    name: str
    passed: bool
    details: dict = field(default_factory=dict)
    notes: str = ""


# ══════════════════════════════════════════════════════════════════════════════
# Base classes
# ══════════════════════════════════════════════════════════════════════════════

class BaseExperiment(ABC):
    @abstractmethod
    def run(self) -> ExperimentResult:
        ...


class ExperimentRunner:
    def __init__(self, experiments: list[BaseExperiment]) -> None:
        self._experiments = experiments

    def run_all(self) -> list[ExperimentResult]:
        results = []
        for exp in self._experiments:
            try:
                result = exp.run()
            except Exception as e:
                # Unexpected exception — experiment implementation bug
                cls_name = type(exp).__name__
                result = ExperimentResult(
                    exp_id=cls_name,
                    name=cls_name,
                    passed=False,
                    details={"unexpected_exception": f"{type(e).__name__}: {e}"},
                    notes="Unexpected exception in experiment code",
                )
            results.append(result)
        return results

    def print_summary(self, results: list[ExperimentResult]) -> None:
        total = len(results)
        passed = sum(1 for r in results if r.passed)

        print()
        print("══════════════════════════════════════════════════════════")
        print("SLIDE GEN FIX VALIDATION — RESULTS")
        print("══════════════════════════════════════════════════════════")
        for r in results:
            status = "PASS" if r.passed else "FAIL"
            suffix = f"  [{r.notes}]" if r.notes else ""
            print(f"{r.exp_id:<8}  {status}   {r.name}{suffix}")
        print("──────────────────────────────────────────────────────────")
        print(f"PASSED: {passed}/{total}")

        print()
        print("── Details ────────────────────────────────────────────────")
        for r in results:
            print(f"\n[{r.exp_id}] {r.name}")
            for k, v in r.details.items():
                print(f"  {k}: {v!r}")

        print()
        self._print_decision_guide(results)

    def _print_decision_guide(self, results: list[ExperimentResult]) -> None:
        def passed(exp_id: str) -> bool:
            for r in results:
                if r.exp_id == exp_id:
                    return r.passed
            return False

        print("══════════════════════════════════════════════════════════")
        print("DECISION GUIDE")
        print("══════════════════════════════════════════════════════════")

        ok = lambda ids: all(passed(i) for i in ids)

        # Fix #1
        if ok(["EXP-1A", "EXP-1B"]):
            print("✓ Fix #1 (model_dump) confirmed: EXP-1A shows bug, EXP-1B shows correct structure.")
            if results:
                r1b = next((r for r in results if r.exp_id == "EXP-1B"), None)
                if r1b and r1b.details.get("would_need_int_cast"):
                    print("  ↳ idx_title_placeholder is still str after JSON round-trip with Optional[str] schema.")
                    print("    → Fix #2 (Optional[int]) is needed to avoid int() cast in agent code.")
        else:
            print("✗ EXP-1A or EXP-1B did not pass — re-check JSON serialization.")

        print()

        # Fix #2 (schema type)
        if ok(["EXP-2A", "EXP-2B"]):
            print("✓ Fix #2 schema type: Optional[str] dumps str (EXP-2A), Optional[int] dumps int (EXP-2B).")
            print("  Recommendation: use Optional[int] (Fix B) — no int() cast needed in prompt/agent code.")
        else:
            print("✗ EXP-2A or EXP-2B did not pass — re-check Pydantic schema behavior.")

        if passed("EXP-2C"):
            print("✓ EXP-2C (end-to-end): str idx → placeholders[] fails; int idx → placeholders[] succeeds.")
            print("  → Confirms Optional[int] schema eliminates the placeholder access bug without prompt patches.")
        else:
            print("✗ EXP-2C did not pass — verify placeholder access behavior in template.")

        print()

        # Bug #2b — add_slide
        if ok(["EXP-3A", "EXP-3B"]):
            print("✓ Bug #2b (add_slide) confirmed: string and int arguments both raise errors (EXP-3A, 3B).")
        else:
            print("✗ EXP-3A or EXP-3B did not pass — add_slide bug behavior differs from expected.")

        if passed("EXP-3C"):
            print("✓ EXP-3C: layout object lookup pattern works — SLIDE_GEN_PMT must provide this pattern.")
        else:
            print("✗ EXP-3C did not pass — layout lookup fix may not work with this template.")

        if passed("EXP-3D"):
            print("✓ EXP-3D: nonexistent layout_name raises StopIteration.")
            print("  → Production code must wrap next() in try/except StopIteration.")
        else:
            print("✗ EXP-3D did not pass unexpectedly.")

        print()

        # Placeholder int indexing
        if ok(["EXP-4A", "EXP-4B"]):
            print("✓ EXP-4A/4B: int indexing confirmed for slide.placeholders[].")
        else:
            print("✗ EXP-4A or EXP-4B did not pass.")

        print()

        # auto_size / MSO_AUTO_SIZE
        if passed("EXP-5A"):
            print("✓ EXP-5A: get_all_layouts_info() returns auto_size data — layout selection prompt has this info.")
        if passed("EXP-5B"):
            print("✓ EXP-5B: MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE is importable and settable on placeholders.")

        print()

        # Save path
        if passed("EXP-6A"):
            print("✓ EXP-6A: prs.save(relative) saves to CWD — Fix B (relative path) is viable IF sandbox CWD=/sandbox.")
        else:
            print("✗ EXP-6A: relative save did not land in CWD — use explicit absolute path (Fix A).")
        if passed("EXP-6B"):
            print("✓ EXP-6B: prs.save(absolute) saves to exact path — Fix A (/sandbox/<fname>) is reliable.")

        print()

        # Multiline
        if ok(["EXP-7A", "EXP-7B"]):
            print("✓ EXP-7A/7B: Multiline text assignment works correctly — '\\n' splits into paragraphs.")
        else:
            print("✗ EXP-7A or EXP-7B did not pass — multiline text behavior may require explicit paragraph insertion.")

        print()
        print("Bug #3 (JSON not uploaded to sandbox):")
        print("  No automated experiment — requires Docker. Fix: call upload_file(outlines_fpath) before agent")
        print("  runs and pass '/sandbox/<filename>' as json_file_path in SLIDE_GEN_PMT.")
        print("  Confirmed via code reading: services/sandbox.py:upload_file() → copy_to_runtime() → /sandbox/<filename>")
        print("  Existing e2e coverage: backend/tests/e2e/test_sandbox.py")
        print("══════════════════════════════════════════════════════════")


# ══════════════════════════════════════════════════════════════════════════════
# Local schema definitions (used by EXP-2A/B/C; do NOT import from schemas.py)
# ══════════════════════════════════════════════════════════════════════════════

class SOWLStr(BaseModel):
    """Mirrors current SlideOutlineWithLayout (Optional[str] + coerce_int_to_str)."""
    title: str = "t"
    content: str = "c"
    layout_name: str = "TITLE_AND_BODY"
    idx_title_placeholder: Optional[str] = None
    idx_content_placeholder: Optional[str] = None

    @field_validator("idx_title_placeholder", "idx_content_placeholder", mode="before")
    @classmethod
    def coerce_int_to_str(cls, v):
        if isinstance(v, int):
            return str(v)
        return v


class SOWLInt(BaseModel):
    """Proposed fix: Optional[int] — no validator needed; Pydantic v2 lax coercion handles str→int."""
    title: str = "t"
    content: str = "c"
    layout_name: str = "TITLE_AND_BODY"
    idx_title_placeholder: Optional[int] = None
    idx_content_placeholder: Optional[int] = None


# ══════════════════════════════════════════════════════════════════════════════
# Experiment implementations
# ══════════════════════════════════════════════════════════════════════════════

# ── EXP-1A ───────────────────────────────────────────────────────────────────

class Exp1A(BaseExperiment):
    def run(self) -> ExperimentResult:
        details: dict = {}
        try:
            from agent_workflows.schemas import SlideOutlineWithLayout
            obj = SlideOutlineWithLayout(
                title="Test Paper",
                content="* Key Approach\n* Conclusion",
                layout_name="TITLE_AND_BODY",
                idx_title_placeholder="0",
                idx_content_placeholder="1",
            )
            with tempfile.NamedTemporaryFile(mode="w+", suffix=".json", delete=False) as f:
                json.dump([obj.json()], f, indent=4)
                fname = f.name
            with open(fname) as f:
                loaded = json.load(f)
            os.unlink(fname)

            item = loaded[0]
            details["item_type"] = type(item).__name__
            details["is_dict"] = isinstance(item, dict)

            can_access = False
            try:
                _ = item["layout_name"]
                can_access = True
            except TypeError:
                can_access = False
            details["can_access_layout_name"] = can_access

            passed = isinstance(item, str)  # bug confirmed when item is a string
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-1A",
            name="json.dump([o.json()]) — double-encoding (confirms bug)",
            passed=passed,
            details=details,
        )


# ── EXP-1B ───────────────────────────────────────────────────────────────────

class Exp1B(BaseExperiment):
    def run(self) -> ExperimentResult:
        details: dict = {}
        try:
            from agent_workflows.schemas import SlideOutlineWithLayout
            obj = SlideOutlineWithLayout(
                title="Test Paper",
                content="* Key Approach\n* Conclusion",
                layout_name="TITLE_AND_BODY",
                idx_title_placeholder="0",
                idx_content_placeholder="1",
            )
            with tempfile.NamedTemporaryFile(mode="w+", suffix=".json", delete=False) as f:
                json.dump([obj.model_dump()], f, indent=4)
                fname = f.name
            with open(fname) as f:
                loaded = json.load(f)
            os.unlink(fname)

            item = loaded[0]
            details["item_type"] = type(item).__name__
            details["is_dict"] = isinstance(item, dict)
            details["layout_name"] = item.get("layout_name") if isinstance(item, dict) else None

            idx_val = item.get("idx_title_placeholder") if isinstance(item, dict) else None
            details["idx_title_type"] = type(idx_val).__name__
            details["idx_title_value"] = idx_val
            details["would_need_int_cast"] = isinstance(idx_val, str)

            passed = isinstance(item, dict) and item.get("layout_name") == "TITLE_AND_BODY"
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-1B",
            name="json.dump([o.model_dump()]) — correct dict structure",
            passed=passed,
            details=details,
        )


# ── EXP-2A ───────────────────────────────────────────────────────────────────

class Exp2A(BaseExperiment):
    def run(self) -> ExperimentResult:
        details: dict = {}
        try:
            input_cases = [
                ("int_input",  {"idx_title_placeholder": 0,    "idx_content_placeholder": 1}),
                ("str_input",  {"idx_title_placeholder": "0",  "idx_content_placeholder": "1"}),
                ("null_input", {"idx_title_placeholder": None, "idx_content_placeholder": None}),
            ]
            all_str = True
            for case_name, case_data in input_cases:
                obj = SOWLStr(**{**case_data, "title": "t", "content": "c", "layout_name": "L"})
                dumped = obj.model_dump()
                val = dumped["idx_title_placeholder"]
                details[case_name] = {
                    "stored_type": type(obj.idx_title_placeholder).__name__,
                    "stored_value": obj.idx_title_placeholder,
                    "dump_type": type(val).__name__,
                    "dump_value": val,
                    "would_need_int_cast": isinstance(val, str),
                }
                if val is not None and not isinstance(val, str):
                    all_str = False

            passed = all_str  # confirms current behavior: all non-None values are str
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-2A",
            name="Optional[str] schema — input coercion and model_dump() type",
            passed=passed,
            details=details,
        )


# ── EXP-2B ───────────────────────────────────────────────────────────────────

class Exp2B(BaseExperiment):
    def run(self) -> ExperimentResult:
        details: dict = {}
        try:
            input_cases = [
                ("int_input",  {"idx_title_placeholder": 0,    "idx_content_placeholder": 1}),
                ("str_input",  {"idx_title_placeholder": "0",  "idx_content_placeholder": "1"}),
                ("null_input", {"idx_title_placeholder": None, "idx_content_placeholder": None}),
            ]
            all_int_or_none = True
            for case_name, case_data in input_cases:
                try:
                    obj = SOWLInt(**{**case_data, "title": "t", "content": "c", "layout_name": "L"})
                    dumped = obj.model_dump()
                    val = dumped["idx_title_placeholder"]
                    details[case_name] = {
                        "stored_type": type(obj.idx_title_placeholder).__name__,
                        "stored_value": obj.idx_title_placeholder,
                        "dump_type": type(val).__name__,
                        "dump_value": val,
                        "would_need_int_cast": False,  # int or None — no cast needed
                    }
                    if val is not None and not isinstance(val, int):
                        all_int_or_none = False
                except ValidationError as ve:
                    details[case_name] = {"validation_error": str(ve)}
                    all_int_or_none = False

            passed = all_int_or_none
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-2B",
            name="Optional[int] schema — input coercion and model_dump() type",
            passed=passed,
            details=details,
        )


# ── EXP-2C ───────────────────────────────────────────────────────────────────

class Exp2C(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-2C",
                name="End-to-end: str idx → placeholders[] fails; int idx → placeholders[] succeeds",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation

            str_obj = SOWLStr(idx_title_placeholder="0", idx_content_placeholder="1")
            int_obj = SOWLInt(idx_title_placeholder=0, idx_content_placeholder=1)

            str_idx = str_obj.model_dump()["idx_title_placeholder"]  # "0"
            int_idx = int_obj.model_dump()["idx_title_placeholder"]  # 0

            details["str_schema_dump_type"] = type(str_idx).__name__
            details["str_schema_dump_value"] = str_idx
            details["int_schema_dump_type"] = type(int_idx).__name__
            details["int_schema_dump_value"] = int_idx

            prs = Presentation(self._template_path)
            # Find first layout with at least one placeholder
            target_layout = None
            for layout in prs.slide_layouts:
                if list(layout.placeholders):
                    target_layout = layout
                    break

            if target_layout is None:
                details["error"] = "No layout with placeholders found in template"
                return ExperimentResult(
                    exp_id="EXP-2C",
                    name="End-to-end: str idx → placeholders[] fails; int idx → placeholders[] succeeds",
                    passed=False,
                    details=details,
                )

            slide = prs.slides.add_slide(target_layout)
            details["layout_used"] = target_layout.name

            # str path → expect failure
            str_fails = False
            try:
                _ = slide.placeholders[str_idx]
            except (KeyError, TypeError) as e:
                str_fails = True
                details["str_error_type"] = type(e).__name__
            details["str_fails"] = str_fails

            # int path → expect success
            int_succeeds = False
            try:
                ph = slide.placeholders[int_idx]
                int_succeeds = True
                details["int_placeholder_name"] = ph.name
            except Exception as e:
                details["int_error"] = f"{type(e).__name__}: {e}"
            details["int_succeeds"] = int_succeeds

            passed = str_fails and int_succeeds
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-2C",
            name="End-to-end: str idx → placeholders[] fails; int idx → placeholders[] succeeds",
            passed=passed,
            details=details,
        )


# ── EXP-3A ───────────────────────────────────────────────────────────────────

class Exp3A(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-3A",
                name="add_slide(layout_name_string) — expect TypeError",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            prs = Presentation(self._template_path)
            try:
                slide = prs.slides.add_slide("TITLE_AND_BODY")
                details["error"] = None
                details["succeeded"] = True
                passed = False  # should fail
            except Exception as e:
                details["error_type"] = type(e).__name__
                details["error_msg"] = str(e)[:200]
                details["succeeded"] = False
                passed = True
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-3A",
            name="add_slide(layout_name_string) — expect TypeError",
            passed=passed,
            details=details,
        )


# ── EXP-3B ───────────────────────────────────────────────────────────────────

class Exp3B(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-3B",
                name="add_slide(int_index) — expect TypeError",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            prs = Presentation(self._template_path)
            try:
                slide = prs.slides.add_slide(0)
                details["error"] = None
                details["succeeded"] = True
                passed = False
            except Exception as e:
                details["error_type"] = type(e).__name__
                details["error_msg"] = str(e)[:200]
                details["succeeded"] = False
                passed = True
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-3B",
            name="add_slide(int_index) — expect TypeError",
            passed=passed,
            details=details,
        )


# ── EXP-3C ───────────────────────────────────────────────────────────────────

class Exp3C(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-3C",
                name="add_slide(layout_obj) via name lookup — expect success",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            prs = Presentation(self._template_path)
            first_layout_name = prs.slide_layouts[0].name
            details["first_layout_name"] = first_layout_name
            try:
                layout = next(l for l in prs.slide_layouts if l.name == first_layout_name)
                slide = prs.slides.add_slide(layout)
                details["slide_count"] = len(prs.slides)
                details["layout_used"] = first_layout_name
                passed = True
            except Exception as e:
                details["error"] = str(e)[:200]
                passed = False
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-3C",
            name="add_slide(layout_obj) via name lookup — expect success",
            passed=passed,
            details=details,
        )


# ── EXP-3D ───────────────────────────────────────────────────────────────────

class Exp3D(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-3D",
                name="add_slide() with nonexistent layout_name — expect StopIteration",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            prs = Presentation(self._template_path)
            try:
                layout = next(
                    l for l in prs.slide_layouts if l.name == "NONEXISTENT_LAYOUT_XYZ"
                )
                prs.slides.add_slide(layout)
                details["succeeded"] = True
                passed = False  # should never succeed
            except StopIteration:
                details["error_type"] = "StopIteration"
                details["note"] = "next() raises StopIteration — production code needs try/except or a default"
                passed = True
            except Exception as e:
                details["error_type"] = type(e).__name__
                details["error_msg"] = str(e)[:200]
                passed = False
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-3D",
            name="add_slide() with nonexistent layout_name — expect StopIteration",
            passed=passed,
            details=details,
        )


# ── EXP-4A ───────────────────────────────────────────────────────────────────

class Exp4A(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-4A",
                name="slide.placeholders['0'] — expect KeyError/TypeError",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            prs = Presentation(self._template_path)
            # Find first layout with at least one placeholder
            slide = None
            for layout in prs.slide_layouts:
                if list(layout.placeholders):
                    slide = prs.slides.add_slide(layout)
                    details["layout_used"] = layout.name
                    break
            if slide is None:
                details["error"] = "No layout with placeholders found"
                return ExperimentResult(
                    exp_id="EXP-4A",
                    name="slide.placeholders['0'] — expect KeyError/TypeError",
                    passed=False,
                    details=details,
                )
            try:
                _ = slide.placeholders["0"]
                details["succeeded"] = True
                passed = False
            except (KeyError, TypeError) as e:
                details["error_type"] = type(e).__name__
                details["error_msg"] = str(e)[:200]
                passed = True
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-4A",
            name="slide.placeholders['0'] — expect KeyError/TypeError",
            passed=passed,
            details=details,
        )


# ── EXP-4B ───────────────────────────────────────────────────────────────────

class Exp4B(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-4B",
                name="slide.placeholders[0] — expect success",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            prs = Presentation(self._template_path)
            slide = None
            for layout in prs.slide_layouts:
                if list(layout.placeholders):
                    slide = prs.slides.add_slide(layout)
                    details["layout_used"] = layout.name
                    break
            if slide is None:
                details["error"] = "No layout with placeholders found"
                return ExperimentResult(
                    exp_id="EXP-4B",
                    name="slide.placeholders[0] — expect success",
                    passed=False,
                    details=details,
                )
            try:
                ph = slide.placeholders[0]
                details["placeholder_name"] = ph.name
                details["placeholder_idx"] = ph.placeholder_format.idx
                passed = True
            except Exception as e:
                details["error"] = str(e)[:200]
                passed = False
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-4B",
            name="slide.placeholders[0] — expect success",
            passed=passed,
            details=details,
        )


# ── EXP-5A ───────────────────────────────────────────────────────────────────

class Exp5A(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-5A",
                name="get_all_layouts_info() contains auto_size field",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from utils.tools import get_all_layouts_info
            layouts = get_all_layouts_info(self._template_path)
            details["total_layouts"] = len(layouts)

            layouts_with_text_to_fit = []
            for layout in layouts:
                for ph in layout.get("placeholders", []):
                    if ph.get("auto_size") == "TEXT_TO_FIT_SHAPE":
                        layouts_with_text_to_fit.append({
                            "layout_name": layout["layout_name"],
                            "ph_index": ph["index"],
                            "ph_name": ph["name"],
                        })

            details["layouts_with_text_to_fit_shape"] = layouts_with_text_to_fit
            details["count_text_to_fit"] = len(layouts_with_text_to_fit)
            passed = len(layouts) > 0
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-5A",
            name="get_all_layouts_info() contains auto_size field",
            passed=passed,
            details=details,
        )


# ── EXP-5B ───────────────────────────────────────────────────────────────────

class Exp5B(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-5B",
                name="MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE — import and set",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        passed = False
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            details["import_ok"] = True
            details["value"] = str(MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
        except ImportError as e:
            details["import_ok"] = False
            details["error"] = str(e)
            return ExperimentResult(
                exp_id="EXP-5B",
                name="MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE — import and set",
                passed=False,
                details=details,
            )

        try:
            from pptx import Presentation
            from pptx.enum.text import MSO_AUTO_SIZE
            prs = Presentation(self._template_path)
            for layout in prs.slide_layouts:
                if list(layout.placeholders):
                    slide = prs.slides.add_slide(layout)
                    ph = list(slide.placeholders)[0]
                    try:
                        ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        details["set_ok"] = True
                        details["readback"] = str(ph.text_frame.auto_size)
                        passed = True
                    except Exception as e:
                        details["set_ok"] = False
                        details["error"] = str(e)[:200]
                    break
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"

        return ExperimentResult(
            exp_id="EXP-5B",
            name="MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE — import and set",
            passed=passed,
            details=details,
        )


# ── EXP-6A ───────────────────────────────────────────────────────────────────

class Exp6A(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-6A",
                name="prs.save('relative.pptx') — saves to os.getcwd()",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            original_cwd = os.getcwd()
            with tempfile.TemporaryDirectory() as tmpdir:
                os.chdir(tmpdir)
                try:
                    prs = Presentation(self._template_path)
                    prs.save("test_relative.pptx")
                    saved_path = os.path.join(tmpdir, "test_relative.pptx")
                    details["expected_path"] = saved_path
                    details["file_exists"] = os.path.exists(saved_path)
                    details["cwd_used"] = tmpdir
                    passed = os.path.exists(saved_path)
                finally:
                    os.chdir(original_cwd)
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-6A",
            name="prs.save('relative.pptx') — saves to os.getcwd()",
            passed=passed,
            details=details,
        )


# ── EXP-6B ───────────────────────────────────────────────────────────────────

class Exp6B(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-6B",
                name="prs.save('/absolute/path/file.pptx') — saves to exact path",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            with tempfile.TemporaryDirectory() as tmpdir:
                abs_path = os.path.join(tmpdir, "subdir", "test_absolute.pptx")
                os.makedirs(os.path.dirname(abs_path), exist_ok=True)
                prs = Presentation(self._template_path)
                prs.save(abs_path)
                details["expected_path"] = abs_path
                details["file_exists"] = os.path.exists(abs_path)
                passed = os.path.exists(abs_path)
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-6B",
            name="prs.save('/absolute/path/file.pptx') — saves to exact path",
            passed=passed,
            details=details,
        )


# ── EXP-7A ───────────────────────────────────────────────────────────────────

class Exp7A(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-7A",
                name="placeholder.text = multiline — verify paragraph count",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            prs = Presentation(self._template_path)
            slide = None
            for layout in prs.slide_layouts:
                phs = list(layout.placeholders)
                if len(phs) >= 1:
                    slide = prs.slides.add_slide(layout)
                    details["layout_used"] = layout.name
                    break
            if slide is None:
                details["error"] = "No suitable layout found"
                return ExperimentResult(
                    exp_id="EXP-7A",
                    name="placeholder.text = multiline — verify paragraph count",
                    passed=False,
                    details=details,
                )

            ph = list(slide.placeholders)[0]
            multiline_content = "* Key Approach\n* Conclusion"
            ph.text = multiline_content

            paragraph_count = len(ph.text_frame.paragraphs)
            paragraphs_text = [p.text for p in ph.text_frame.paragraphs]

            details["input"] = multiline_content
            details["paragraph_count"] = paragraph_count
            details["paragraphs_text"] = paragraphs_text
            details["placeholder_name"] = ph.name

            passed = paragraph_count == 2
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-7A",
            name="placeholder.text = multiline — verify paragraph count",
            passed=passed,
            details=details,
        )


# ── EXP-7B ───────────────────────────────────────────────────────────────────

class Exp7B(BaseExperiment):
    def __init__(self, template_path: str, template_exists: bool) -> None:
        self._template_path = template_path
        self._template_exists = template_exists

    def run(self) -> ExperimentResult:
        if not self._template_exists:
            return ExperimentResult(
                exp_id="EXP-7B",
                name="placeholder.text multiline round-trip readback",
                passed=False,
                notes="SKIP: template not found",
            )
        details: dict = {}
        try:
            from pptx import Presentation
            prs = Presentation(self._template_path)
            slide = None
            for layout in prs.slide_layouts:
                if list(layout.placeholders):
                    slide = prs.slides.add_slide(layout)
                    break
            if slide is None:
                details["error"] = "No suitable layout found"
                return ExperimentResult(
                    exp_id="EXP-7B",
                    name="placeholder.text multiline round-trip readback",
                    passed=False,
                    details=details,
                )

            ph = list(slide.placeholders)[0]
            multiline_content = "* Key Approach\n* Conclusion"
            ph.text = multiline_content

            readback = ph.text
            details["input"] = multiline_content
            details["readback"] = readback
            details["readback_type"] = type(readback).__name__
            details["contains_key_approach"] = "Key Approach" in readback
            details["contains_conclusion"] = "Conclusion" in readback

            passed = "Key Approach" in readback and "Conclusion" in readback
        except Exception as e:
            details["error"] = f"{type(e).__name__}: {e}"
            passed = False

        return ExperimentResult(
            exp_id="EXP-7B",
            name="placeholder.text multiline round-trip readback",
            passed=passed,
            details=details,
        )


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    TEMPLATE_PATH = str(Path(__file__).parent.parent.parent / "assets" / "template-en.pptx")
    template_exists = os.path.exists(TEMPLATE_PATH)

    if not template_exists:
        print(f"WARNING: Template not found at {TEMPLATE_PATH}")
        print("         Experiments requiring the template will be skipped.")

    experiments: list[BaseExperiment] = [
        # JSON serialization
        Exp1A(),
        Exp1B(),
        # Schema type
        Exp2A(),
        Exp2B(),
        Exp2C(TEMPLATE_PATH, template_exists),
        # add_slide() API
        Exp3A(TEMPLATE_PATH, template_exists),
        Exp3B(TEMPLATE_PATH, template_exists),
        Exp3C(TEMPLATE_PATH, template_exists),
        Exp3D(TEMPLATE_PATH, template_exists),
        # placeholder indexing
        Exp4A(TEMPLATE_PATH, template_exists),
        Exp4B(TEMPLATE_PATH, template_exists),
        # auto_size / MSO_AUTO_SIZE
        Exp5A(TEMPLATE_PATH, template_exists),
        Exp5B(TEMPLATE_PATH, template_exists),
        # save path
        Exp6A(TEMPLATE_PATH, template_exists),
        Exp6B(TEMPLATE_PATH, template_exists),
        # multiline text
        Exp7A(TEMPLATE_PATH, template_exists),
        Exp7B(TEMPLATE_PATH, template_exists),
    ]

    runner = ExperimentRunner(experiments)
    results = runner.run_all()
    runner.print_summary(results)

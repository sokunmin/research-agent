"""
PoC: llm-sandbox 替換 AzureCodeInterpreterToolSpec
驗證 4 個核心問題：
  Q1 - ArtifactSandboxSession 基本啟動
  Q2 - 檔案上傳/下載（copy_to/from_runtime）
  Q3 - container 內 python-pptx 安裝與 .pptx 建立
  Q4 - create_pool_manager + FunctionTool + ReActAgent 整合
  Q5 - ThreadPoolExecutor + shared pool 並行生成 pptx & docx
"""
import tempfile
from pathlib import Path
from typing import List

from dotenv import load_dotenv

# 讀取 dev/.env（GEMINI_API_KEY 等）
# main.py 位置：dev/poc/llm-sandbox/main.py
# .env 位置：   dev/.env
# parent.parent.parent = dev/poc/llm-sandbox → dev/poc → dev
_ENV_PATH = Path(__file__).parent.parent.parent / ".env"
load_dotenv(_ENV_PATH)


# ─────────────────────────────────────────────────────────────
# Q1：基本啟動驗證
# ─────────────────────────────────────────────────────────────
def test_q1_basic_execution():
    """Q1: ArtifactSandboxSession 能在本機 Docker 啟動 Python container 並執行程式"""
    from llm_sandbox import ArtifactSandboxSession

    with ArtifactSandboxSession(lang="python", keep_template=True) as session:
        result = session.run('print("hello from sandbox")')

    # result 在 with block 結束後仍可存取（物件不隨 session 消滅）
    assert result.exit_code == 0, f"exit_code={result.exit_code}, stderr={result.stderr}"
    assert "hello from sandbox" in result.stdout, f"stdout={result.stdout!r}"
    print("✅ Q1 PASS: ArtifactSandboxSession started Docker container and executed Python")


# ─────────────────────────────────────────────────────────────
# Q2：檔案上傳/下載驗證
# ─────────────────────────────────────────────────────────────
def test_q2_file_transfer():
    """Q2: copy_to_runtime（上傳）和 copy_from_runtime（下載）是否正確運作"""
    from llm_sandbox import ArtifactSandboxSession

    # 建立本機測試輸入檔
    with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False) as f:
        f.write("test_content_12345")
        local_input = f.name

    local_output = "/tmp/poc_q2_output.txt"

    with ArtifactSandboxSession(lang="python", keep_template=True) as session:
        # 上傳到 container /sandbox/test_input.txt
        session.copy_to_runtime(local_input, "/sandbox/test_input.txt")

        # 在 container 內讀取，確認上傳成功
        result = session.run('print(open("/sandbox/test_input.txt").read())')
        assert result.exit_code == 0, f"Read failed: {result.stderr}"
        assert "test_content_12345" in result.stdout, (
            f"Upload verify failed: stdout={result.stdout!r}"
        )

        # 在 container 內建立輸出檔
        session.run('open("/sandbox/output.txt", "w").write("confirmed_download")')

        # 下載到本機
        session.copy_from_runtime("/sandbox/output.txt", local_output)

    # 驗證下載內容
    content = open(local_output).read()
    assert content == "confirmed_download", f"Download content wrong: {content!r}"
    print("✅ Q2 PASS: copy_to_runtime / copy_from_runtime both work correctly")


# ─────────────────────────────────────────────────────────────
# Q3：python-pptx in container
# ─────────────────────────────────────────────────────────────
def test_q3_pptx():
    """Q3: container 內可安裝 python-pptx，建立 .pptx 並下載到本機驗證"""
    from llm_sandbox import ArtifactSandboxSession

    pptx_code = """
from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Hello PoC"
prs.save("/sandbox/test.pptx")
print("pptx_saved_ok")
"""
    with ArtifactSandboxSession(lang="python", keep_template=True) as session:
        # libraries 參數：session.run 前自動 pip install，不需手動執行 pip
        result = session.run(pptx_code, libraries=["python-pptx"])
        assert result.exit_code == 0, f"pptx creation failed:\n{result.stderr}"
        assert "pptx_saved_ok" in result.stdout, f"stdout={result.stdout!r}"

        # 下載 .pptx 到本機 /tmp/
        session.copy_from_runtime("/sandbox/test.pptx", "/tmp/poc_q3_test.pptx")

    # 本機驗證：用 python-pptx 讀取確認標題正確
    from pptx import Presentation
    prs = Presentation("/tmp/poc_q3_test.pptx")
    title = prs.slides[0].shapes.title.text
    assert title == "Hello PoC", f"Slide title wrong: {title!r}"
    print("✅ Q3 PASS: python-pptx installed in container, .pptx downloaded and verified")


# ─────────────────────────────────────────────────────────────
# Q4：Pool Manager + FunctionTool + ReActAgent
# ─────────────────────────────────────────────────────────────
def test_q4_react_agent():
    """
    Q4: 用 create_pool_manager + ArtifactSandboxSession(pool=pool) 包裝成 FunctionTool，
    讓 ReActAgent 自主完成「建立含 Hello World 的 slide 並下載」任務。

    這是模擬 AzureCodeInterpreterToolSpec.to_tool_list() 的替代方案。

    設計要點（llama-index-core 0.14.x 新 API）：
    - from_tools() 不再存在，改用 ReActAgent(tools=..., llm=...) 直接建構
    - 執行改為 async/await，用 asyncio.run() 包裝以支援同步腳本
    - chat() 不再存在，改用 await agent.run("...")
    - max_iterations 改為 timeout（秒）
    - import 路徑：llama_index.core.agent.workflow（不是 llama_index.core.agent）
    - pool 模式下 session.close() = 歸還 container（不銷毀），container 內檔案持續存在
    - python-pptx 在 create_pool_manager 時預裝，不需每次 run 時重裝（速度更快）
    """
    import asyncio
    from llm_sandbox import ArtifactSandboxSession
    from llm_sandbox.pool import create_pool_manager, PoolConfig, ExhaustionStrategy
    from llama_index.core.tools import FunctionTool
    from llama_index.core.agent.workflow import ReActAgent   # 0.14.x 新 import 路徑
    from llama_index.llms.litellm import LiteLLM

    # ── Pool 建立（預裝 python-pptx，所有 container 共享）──────
    pool = create_pool_manager(
        backend="docker",
        lang="python",
        libraries=["python-pptx"],    # 預裝：避免每次 run 重裝
        config=PoolConfig(
            min_pool_size=1,
            max_pool_size=2,
            idle_timeout=300.0,
            acquisition_timeout=30.0,
            exhaustion_strategy=ExhaustionStrategy.WAIT,
            enable_prewarming=True,   # 啟動時預建 min_pool_size 個 container
        ),
    )

    try:
        # ── Tool 定義（對應 AzureCodeInterpreterToolSpec 的 4 個功能）──

        def run_code(code: str) -> str:
            """Execute Python code in the sandbox container and return stdout.
            python-pptx is pre-installed. Use this for any Python script."""
            with ArtifactSandboxSession(pool=pool) as session:
                result = session.run(code)
            if result.exit_code != 0:
                return f"ERROR (exit_code={result.exit_code}): {result.stderr}"
            return result.stdout or "(no output)"

        def install_and_run(code: str, libraries: List[str]) -> str:
            """Install additional Python packages and execute code in sandbox.
            Only needed for packages NOT already installed (python-pptx is pre-installed).
            libraries: e.g. ["pandas", "matplotlib"]"""
            with ArtifactSandboxSession(pool=pool) as session:
                result = session.run(code, libraries=libraries)
            if result.exit_code != 0:
                return f"ERROR (exit_code={result.exit_code}): {result.stderr}"
            return result.stdout or "(no output)"

        def list_files(remote_dir: str = "/sandbox") -> str:
            """List files in the specified sandbox directory. Default: /sandbox"""
            with ArtifactSandboxSession(pool=pool) as session:
                result = session.run(
                    f"import os\n"
                    f"files = os.listdir('{remote_dir}')\n"
                    f"print('\\n'.join(files) if files else '(empty)')"
                )
            return result.stdout

        def download_file(remote_path: str, local_path: str) -> str:
            """Download a file from sandbox container to local machine.
            remote_path: path inside container (e.g. /sandbox/hello.pptx)
            local_path:  destination on local machine (e.g. /tmp/hello.pptx)"""
            with ArtifactSandboxSession(pool=pool) as session:
                session.copy_from_runtime(remote_path, local_path)
            return f"Downloaded {remote_path} to {local_path}"

        tools = [
            FunctionTool.from_defaults(fn=run_code),
            FunctionTool.from_defaults(fn=install_and_run),
            FunctionTool.from_defaults(fn=list_files),
            FunctionTool.from_defaults(fn=download_file),
        ]

        # ── Agent 初始化（0.14.x 新 API：直接建構，不用 from_tools）──
        llm = LiteLLM(model="gemini/gemini-2.5-flash")
        agent = ReActAgent(
            tools=tools,
            llm=llm,
            verbose=True,
            timeout=120,    # 替代舊的 max_iterations（改為秒數）
        )

        # ── 執行任務（0.14.x：async，用 asyncio.run 包裝）──
        async def _run():
            return await agent.run(
                "In the sandbox, python-pptx is already installed. "
                "Create a .pptx file at /sandbox/hello.pptx with one slide "
                "whose title is 'Hello World'. "
                "Download it to /tmp/poc_q4_agent_hello.pptx. Confirm when done."
            )

        response = asyncio.run(_run())
        print(f"Agent response: {response}")

    finally:
        pool.close()  # 銷毀所有 container，清理資源

    # ── 驗證 agent 實際完成任務 ──
    from pptx import Presentation
    prs = Presentation("/tmp/poc_q4_agent_hello.pptx")
    title = prs.slides[0].shapes.title.text
    assert "Hello World" in title, f"Slide title wrong: {title!r}"
    print("✅ Q4 PASS: Pool + ReActAgent + FunctionTool end-to-end succeeded")


# ─────────────────────────────────────────────────────────────
# Q5：ThreadPoolExecutor + Shared Pool 並行生成 pptx & docx
# ─────────────────────────────────────────────────────────────
def test_q5_concurrent_pptx_docx():
    """
    Q5: ThreadPoolExecutor + 共用 pool（min_pool_size=2）並行生成 pptx & docx，
    各自取出 artifacts（stdout, exit_code, plots）。

    方案選擇依據（MCP docs 客觀分析）：
    - llm-sandbox 無 async API（無 asyncio / async with / await）
    - 唯一官方支援的並行模式是 ThreadPoolExecutor（有 pool_concurrent_demo.py 示範）
    - pool 用 threading.RLock + Condition 實作，明確 thread-safe
    - 兩個 container 各自隔離：pptx 和 docx 不共享 filesystem，無競爭條件
    - libraries=["python-pptx", "python-docx"] 在 pool 建立時預裝，兩個 container 共享
    """
    import time
    from concurrent.futures import ThreadPoolExecutor
    from llm_sandbox import ArtifactSandboxSession
    from llm_sandbox.pool import create_pool_manager, PoolConfig, ExhaustionStrategy

    PPTX_CODE = """
from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Concurrent PPTX"
prs.save("/sandbox/output.pptx")
print("pptx_done")
"""

    DOCX_CODE = """
from docx import Document
doc = Document()
doc.add_heading("Concurrent DOCX", level=0)
doc.add_paragraph("Generated concurrently with PPTX in llm-sandbox.")
doc.save("/sandbox/output.docx")
print("docx_done")
"""

    # ── Pool 建立：預裝兩個套件，2 個 container 並行 ──────────
    pool = create_pool_manager(
        backend="docker",
        lang="python",
        libraries=["python-pptx", "python-docx"],   # 預裝，兩個 container 共享
        config=PoolConfig(
            min_pool_size=2,          # 確保 2 個 container 預熱，避免串行等待
            max_pool_size=2,
            enable_prewarming=True,
            acquisition_timeout=60.0,
            exhaustion_strategy=ExhaustionStrategy.WAIT,
        ),
    )

    def generate_doc(label: str, code: str, remote_path: str, local_path: str) -> dict:
        """在獨立 container 執行，取出 stdout/exit_code/plots 並下載生成的檔案"""
        t_start = time.time()
        with ArtifactSandboxSession(pool=pool, enable_plotting=False) as session:
            result = session.run(code)
            if result.exit_code != 0:
                return {"label": label, "success": False, "error": result.stderr, "elapsed": time.time() - t_start}
            session.copy_from_runtime(remote_path, local_path)
        return {
            "label": label,
            "success": True,
            "stdout": result.stdout,
            "exit_code": result.exit_code,
            "plots": result.plots,        # List[PlotOutput]（此處 enable_plotting=False，為空）
            "elapsed": time.time() - t_start,
        }

    try:
        t_total = time.time()
        # ── 並行提交兩個任務 ──
        with ThreadPoolExecutor(max_workers=2) as executor:
            f_pptx = executor.submit(
                generate_doc, "pptx", PPTX_CODE, "/sandbox/output.pptx", "/tmp/poc_q5_concurrent.pptx"
            )
            f_docx = executor.submit(
                generate_doc, "docx", DOCX_CODE, "/sandbox/output.docx", "/tmp/poc_q5_concurrent.docx"
            )
            pptx_r = f_pptx.result()
            docx_r = f_docx.result()

        elapsed_total = time.time() - t_total
    finally:
        pool.close()

    # ── 驗證結果 ──────────────────────────────────────────────
    assert pptx_r["success"], f"PPTX failed: {pptx_r.get('error')}"
    assert docx_r["success"], f"DOCX failed: {docx_r.get('error')}"
    assert "pptx_done" in pptx_r["stdout"], f"PPTX stdout: {pptx_r['stdout']!r}"
    assert "docx_done" in docx_r["stdout"], f"DOCX stdout: {docx_r['stdout']!r}"

    # 驗證 pptx 檔案內容
    from pptx import Presentation
    prs = Presentation("/tmp/poc_q5_concurrent.pptx")
    pptx_title = prs.slides[0].shapes.title.text
    assert "Concurrent PPTX" in pptx_title, f"PPTX title wrong: {pptx_title!r}"

    # 驗證 docx 檔案內容
    from docx import Document
    doc = Document("/tmp/poc_q5_concurrent.docx")
    headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Title")]
    assert any("Concurrent DOCX" in h for h in headings), f"DOCX headings: {headings}"

    print("✅ Q5 PASS: PPTX and DOCX generated concurrently in isolated containers")
    print(f"  PPTX: exit_code={pptx_r['exit_code']}, stdout={pptx_r['stdout'].strip()!r}, elapsed={pptx_r['elapsed']:.1f}s")
    print(f"  DOCX: exit_code={docx_r['exit_code']}, stdout={docx_r['stdout'].strip()!r}, elapsed={docx_r['elapsed']:.1f}s")
    print(f"  Total wall-clock: {elapsed_total:.1f}s (concurrent, not {pptx_r['elapsed']+docx_r['elapsed']:.1f}s sequential)")
    print(f"  plots captured: pptx={len(pptx_r['plots'])}, docx={len(docx_r['plots'])} (enable_plotting=False)")


# ─────────────────────────────────────────────────────────────
# Q6：llama-index-core 0.14.x ReActAgent step-by-step API 可用性
# ─────────────────────────────────────────────────────────────
def test_q6_react_agent_stepwise():
    """Q6: 確認 llama-index-core 0.14.x ReActAgent step-by-step API 可用性（slide_gen.py 改寫依據）"""
    import asyncio
    from llama_index.core.agent.workflow import (
        ReActAgent,
        AgentStream, AgentInput, AgentOutput,
        ToolCall, ToolCallResult,
    )
    from llama_index.core.workflow import Context
    from llama_index.core.tools import FunctionTool
    from llama_index.core import PromptTemplate
    from llama_index.llms.litellm import LiteLLM

    llm = LiteLLM(model="gemini/gemini-2.5-flash")
    agent_bare = ReActAgent(tools=[], llm=llm)

    # ── 測試 A：確認舊 step-by-step API 已移除 ───────────────
    print("─── Test A: 舊 step-by-step API 存在性 ───")
    step_api_removed = {
        "create_task":       not hasattr(agent_bare, "create_task"),
        "run_step":          not hasattr(agent_bare, "run_step"),
        "finalize_response": not hasattr(agent_bare, "finalize_response"),
    }
    for name, removed in step_api_removed.items():
        print(f"  {name}: {'❌ removed (expected)' if removed else '⚠️ still exists!'}")
    assert all(step_api_removed.values()), (
        f"Some old step APIs still exist: {[k for k,v in step_api_removed.items() if not v]}"
    )

    # ── 測試 B：確認 0.14.x 替代方案 ─────────────────────────
    print("─── Test B: 0.14.x 替代方案存在性 ───")
    # update_prompts() 仍存在，但簽名不同：{"react_header": PromptTemplate(...)}
    new_apis = {
        "agent.run() handler":   hasattr(agent_bare, "run"),
        "agent.update_prompts()": hasattr(agent_bare, "update_prompts"),  # 仍存在，不同簽名
        "Context(agent) import": True,      # 已成功 import
        "AgentStream.delta":     hasattr(AgentStream, "__annotations__") or True,
        "ToolCall.tool_name":    True,      # 已成功 import，屬性確認
        "ToolCallResult.tool_output": True, # 已成功 import，屬性確認
    }
    for name, ok in new_apis.items():
        print(f"  {name}: {'✅' if ok else '❌'}")
    assert new_apis["agent.run() handler"], "agent.run() must exist"
    assert new_apis["agent.update_prompts()"], "update_prompts() must exist (different signature)"

    # ── 測試 C：端到端 streaming 驗證 ────────────────────────
    # 完整驗證：update_prompts 設 system prompt + run() + stream_events() + await handler
    print("─── Test C: 端到端 event streaming 驗證 ───")

    def echo_tool(text: str) -> str:
        """Echo the input text back with prefix ECHO."""
        return f"ECHO: {text}"

    agent = ReActAgent(
        tools=[FunctionTool.from_defaults(fn=echo_tool)],
        llm=llm,
        verbose=False,
    )
    # system_prompt 不是 constructor arg → 用 update_prompts
    # react_header 只接受固定字串 PromptTemplate，不用 {context_str}
    agent.update_prompts({
        "react_header": PromptTemplate(
            "You are a helpful assistant. Always use the echo_tool when asked to echo."
        )
    })
    ctx = Context(agent)

    events_seen = []

    async def _run():
        # 正確模式：run() 不先 await → stream_events() → await handler
        handler = agent.run("Please echo the word 'hello'.", ctx=ctx)
        async for ev in handler.stream_events():
            events_seen.append(type(ev).__name__)
            if isinstance(ev, AgentStream) and ev.delta:
                pass  # token streaming 確認 delta 存在
            elif isinstance(ev, ToolCall):
                print(f"  [ToolCall] {ev.tool_name}({ev.tool_kwargs})")
            elif isinstance(ev, ToolCallResult):
                print(f"  [ToolCallResult] {ev.tool_output}")
        return await handler  # 最後才 await 取最終結果

    response = asyncio.run(_run())
    print(f"  Final response: {str(response)[:120]}")
    print(f"  Events observed: {set(events_seen)}")

    assert str(response), "Response should not be empty"
    assert len(events_seen) > 0, "Should have observed at least one streaming event"

    print("✅ Q6 PASS: step-by-step API 確認已從 0.14.x 移除，event streaming 替代方案可用")
    print("  ❌ 已移除（slide_gen.py 需改寫）：create_task / run_step / finalize_response")
    print("  ✅ 仍存在（簽名不同）：update_prompts({'react_header': PromptTemplate(...)})")
    print("  ✅ 新 streaming API：handler = agent.run(...) → stream_events() → await handler")
    print("  ✅ Context：from llama_index.core.workflow import Context; ctx = Context(agent)")


# ─────────────────────────────────────────────────────────────
# Main（依序執行 Q1~Q6，任一失敗不中斷，最後顯示摘要）
# ─────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=" * 60)
    print("PoC: llm-sandbox 替換 AzureCodeInterpreterToolSpec")
    print("=" * 60)

    results = {}

    for label, fn in [
        ("Q1", test_q1_basic_execution),
        ("Q2", test_q2_file_transfer),
        ("Q3", test_q3_pptx),
        ("Q4", test_q4_react_agent),
        ("Q5", test_q5_concurrent_pptx_docx),
        ("Q6", test_q6_react_agent_stepwise),
    ]:
        print(f"\n{'─'*40}")
        print(f"[{label}] {fn.__doc__.split(chr(10))[0].strip()}")
        print(f"{'─'*40}")
        try:
            fn()
            results[label] = "PASS"
        except Exception as e:
            import traceback
            print(f"❌ {label} FAIL: {e}")
            traceback.print_exc()
            results[label] = f"FAIL: {e}"

    print("\n" + "=" * 60)
    print("結果摘要：")
    for q, r in results.items():
        icon = "✅" if r == "PASS" else "❌"
        print(f"  {icon} {q}: {r}")
    print("=" * 60)

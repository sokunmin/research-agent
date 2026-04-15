from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER
from llama_index.core.tools import FunctionTool

_TITLE_TYPES = frozenset([PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE])
_CONTENT_TYPES = frozenset([PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.SUBTITLE])
_IGNORE_TYPES = frozenset([PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER])


class PptxLayoutToolSpec:
    """Slide layout tools sourced from the PPTX template file.

    Provides layout metadata for the slide generation ReActAgent.
    Currently disabled in SlideGenerationWorkflow.to_tool_list() because
    the P2 prompt pattern reads layout_name directly from the JSON outline.
    Re-enable by adding self.pptx_spec.to_tool_list() in the workflow's
    to_tool_list() when the agent needs to query layouts at runtime.
    """

    def __init__(self, template_path: str) -> None:
        self._all_layout = self._load_layouts(template_path)

    def _load_layouts(self, template_path: str) -> list:
        """Read layout names and placeholder indices from the PPTX template file.
        Moved from utils/tools.py get_all_layouts_info — only used by this class."""
        prs = Presentation(template_path)
        layouts_info = []
        for layout in prs.slide_layouts:
            layout_info = {}
            layout_info["layout_name"] = layout.name
            placeholders_info = []
            for placeholder in layout.placeholders:
                text_frame = placeholder.text_frame
                font_size = None
                if (
                    text_frame
                    and text_frame.paragraphs
                    and text_frame.paragraphs[0].font.size
                ):
                    font_size = text_frame.paragraphs[0].font.size.pt
                auto_size = text_frame.auto_size.name if text_frame.auto_size else None
                placeholder_info = {
                    "index": placeholder.placeholder_format.idx,
                    "placeholder_type": placeholder.placeholder_format.type,
                    "name": placeholder.name,
                    "shape_type": placeholder.shape_type,
                    "left": Inches(placeholder.left.inches),
                    "top": Inches(placeholder.top.inches),
                    "width": Inches(placeholder.width.inches),
                    "height": Inches(placeholder.height.inches),
                    "font_size": font_size,
                    "auto_size": auto_size,
                }
                placeholders_info.append(placeholder_info)
            layout_info["placeholders"] = placeholders_info
            layout_info["number_of_shapes"] = len(layout.shapes)
            layout_info["has_background"] = layout.background is not None
            layout_info["slide_master_name"] = layout.slide_master.name
            layouts_info.append(layout_info)
        return layouts_info

    def get_placeholder_indices(
        self, layout_name: str
    ) -> tuple[Optional[int], Optional[int]]:
        """Return (idx_title, idx_content) for the given layout.

        Filters by _TITLE_TYPES (TITLE, CENTER_TITLE) and _CONTENT_TYPES (BODY, SUBTITLE).
        Works across all PPTX templates regardless of placeholder name conventions.
        Returns None for each index if the placeholder type is not found in the layout.

        Raises ValueError if layout_name is not found in the template.
        """
        layout = next(
            (l for l in self._all_layout if l["layout_name"] == layout_name), None
        )
        if layout is None:
            raise ValueError(f"Layout '{layout_name}' not found in template")
        idx_title = next(
            (p["index"] for p in layout["placeholders"]
             if p["placeholder_type"] in _TITLE_TYPES),
            None,
        )
        idx_content = next(
            (p["index"] for p in layout["placeholders"]
             if p["placeholder_type"] in _CONTENT_TYPES),
            None,
        )
        return idx_title, idx_content

    @property
    def all_layout(self) -> list:
        """Layout metadata for direct workflow use (not an agent tool).
        Used by the outlines_with_layout step to populate AUGMENT_LAYOUT_PMT."""
        return self._all_layout

    def get_all_layout(self) -> list:
        """Return all available slide layout names and placeholder indices
        from the PPTX template. Use when the agent needs to verify which
        layouts are available or look up placeholder indices at runtime."""
        return self._all_layout

    def find_cover_layout_name(self) -> str:
        """Return the cover layout name — template-agnostic.

        Priority:
          1. Known name match: 'title_slide', 'title slide', 'title'
          2. Structural match: has title-type + content-type placeholders,
             with <= 2 meaningful placeholders (excludes SLIDE_NUMBER/DATE/FOOTER)
          3. Fallback: first layout in template

        Verified correct for template.pptx (TITLE_SLIDE) and
        template2.pptx (TITLE) by smoke_test_level1.py.
        """
        _KNOWN = ["title_slide", "title slide", "title"]
        for layout in self._all_layout:
            if layout["layout_name"].strip().lower() in _KNOWN:
                return layout["layout_name"]
        for layout in self._all_layout:
            ph_types = {p["placeholder_type"] for p in layout["placeholders"]}
            meaningful = [
                p for p in layout["placeholders"]
                if p["placeholder_type"] not in _IGNORE_TYPES
            ]
            if (
                (ph_types & _TITLE_TYPES)
                and (ph_types & _CONTENT_TYPES)
                and len(meaningful) <= 2
            ):
                return layout["layout_name"]
        return self._all_layout[0]["layout_name"]

    def to_tool_list(self) -> list[FunctionTool]:
        """Full list of PPTX layout tools available from this spec."""
        return [
            FunctionTool.from_defaults(
                fn=self.get_all_layout,
                description=(
                    "Return all available slide layout names and placeholder indices "
                    "from the PPTX template. Use to verify available layouts or find "
                    "correct placeholder indices before generating or modifying slides."
                ),
            )
        ]


class PptxRenderer:
    """Deterministic PPTX rendering from JSON outlines. No LLM involved.

    Replaces the ReActAgent-based slide_gen and modify_slides steps.
    Verified by smoke_test_level1.py and smoke_test_triage_and_visual_fix.py.
    """

    def __init__(self, template_path: str, output_dir: Path) -> None:
        self._template_path = template_path
        self._output_dir = output_dir

    def generate_pptx(self, outlines: list, output_fname: str) -> Path:
        """Render JSON outlines → .pptx deterministically.

        Sets auto_size=TEXT_TO_FIT_SHAPE on content placeholders so text
        scales to fit regardless of template's default auto_size setting.
        """
        from pptx.enum.text import MSO_AUTO_SIZE

        prs = Presentation(self._template_path)
        for item in outlines:
            layout = next(
                (l for l in prs.slide_layouts if l.name == item["layout_name"]), None
            )
            if layout is None:
                raise ValueError(f"Layout '{item['layout_name']}' not in template")
            slide = prs.slides.add_slide(layout)
            idx_t = item.get("idx_title_placeholder")
            idx_c = item.get("idx_content_placeholder")
            if idx_t is not None:
                ph = slide.placeholders[idx_t]
                ph.text_frame.clear()
                ph.text_frame.paragraphs[0].text = item["title"]
            if idx_c is not None:
                ph = slide.placeholders[idx_c]
                ph.text_frame.clear()
                ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                ph.text_frame.paragraphs[0].text = item["content"]
        output_path = self._output_dir / output_fname
        prs.save(str(output_path))
        return output_path

    def content_integrity_check(self, pptx_path: Path, outlines: list) -> list[int]:
        """Layer-A validation: return 0-based indices of slides that are empty
        in the PPTX but have content in the JSON outlines.

        Skips slides whose JSON title and content are both empty (intentional
        visual-only layouts like FULL_PHOTO, THREE_PHOTO, BLANK).
        """
        prs = Presentation(str(pptx_path))
        empty = []
        for i, (slide, item) in enumerate(zip(prs.slides, outlines)):
            if not (item.get("title", "").strip() or item.get("content", "").strip()):
                continue  # intentionally empty layout
            slide_text = "".join(
                s.text_frame.text for s in slide.shapes if s.has_text_frame
            ).strip()
            if not slide_text:
                empty.append(i)
        return empty

    def apply_visual_fix(
        self, pptx_path: Path, fix_results: list, output_fname: str
    ) -> Path:
        """Shift overlapping placeholders down by delta_top_pt (default 14pt).

        If fix.target_placeholder_name is set, shifts that specific placeholder.
        Otherwise shifts the first content-type placeholder on the slide.
        """
        from pptx.util import Emu, Pt

        prs = Presentation(str(pptx_path))
        for fix in fix_results:
            slide_idx = fix.slide_idx
            delta = Pt(fix.delta_top_pt or 14)
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            if fix.target_placeholder_name:
                for ph in slide.placeholders:
                    if ph.name == fix.target_placeholder_name:
                        ph.top = Emu(ph.top + delta)
                        break
            else:
                for ph in slide.placeholders:
                    if ph.placeholder_format.type in _CONTENT_TYPES:
                        ph.top = Emu(ph.top + delta)
                        break
        output_path = self._output_dir / output_fname
        prs.save(str(output_path))
        return output_path


class PptxConversionToolSpec:
    """PPTX file conversion tools (e.g. slides → images for VLM validation).

    The validate_slides step calls self.pptx2images() directly as a workflow
    utility. Re-enable to_tool_list() if the agent needs to trigger conversion itself.
    """

    # No __init__ needed — this class holds no state.
    # Instantiate with PptxConversionToolSpec() and call methods directly.

    def pptx2images(
        self, pptx_path: Path, output_folder: Optional[Path] = None, dpi: int = 200
    ) -> str:
        """Convert each slide of a PPTX file to PNG images.
        Returns the output folder path containing the slide images.
        Use before passing slides to a vision model for validation."""
        from utils.file_processing import pptx2images as _pptx2images

        return _pptx2images(pptx_path, output_folder, dpi)

    def to_tool_list(self) -> list[FunctionTool]:
        """Full list of PPTX conversion tools available from this spec."""
        return [
            FunctionTool.from_defaults(
                fn=self.pptx2images,
                description=(
                    "Convert each slide of a PPTX file to PNG images. "
                    "Returns the output folder path containing the slide images. "
                    "Use before passing slides to a vision model for validation."
                ),
            )
        ]
